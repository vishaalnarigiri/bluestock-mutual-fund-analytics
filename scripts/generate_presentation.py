import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs_path = "/Users/vaishnavnarigiri/Desktop/bluestock/reports/Presentation.pptx"
charts_dir = "/Users/vaishnavnarigiri/Desktop/bluestock/reports/charts"

def create_presentation():
    prs = Presentation()
    
    # Custom color palette (Fintech theme)
    navy_blue = RGBColor(30, 58, 138)    # #1e3a8e
    slate_gray = RGBColor(71, 85, 105)   # #475569
    dark_slate = RGBColor(15, 23, 42)    # #0f172a
    teal_green = RGBColor(15, 118, 110)  # #0f766e
    
    # Helper to apply slide formatting
    def add_slide_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Arial'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = navy_blue
        return title_box

    # Helper to add standard text contents
    def add_bullet_content(slide, bullets, left=Inches(0.5), top=Inches(1.3), width=Inches(5.0), height=Inches(5.0), font_size=16):
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        
        for idx, bullet in enumerate(bullets):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = bullet
            p.font.name = 'Arial'
            p.font.size = Pt(font_size)
            p.font.color.rgb = dark_slate
            # Sub-bullet check
            if bullet.strip().startswith("•"):
                p.level = 1
                p.text = bullet.replace("•", "").strip()
            else:
                p.level = 0
                p.space_after = Pt(8)
        return tx_box

    # --- SLIDE 1: COVER SLIDE ---
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Large Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "MUTUAL FUND ANALYTICS PLATFORM"
    p.font.name = 'Arial'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = navy_blue
    
    p2 = tf.add_paragraph()
    p2.text = "End-to-End ETL, Performance Analytics & Dashboard Solutions"
    p2.font.name = 'Arial'
    p2.font.size = Pt(18)
    p2.font.color.rgb = slate_gray
    p2.space_before = Pt(10)
    
    # Metadata
    meta_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(8.0), Inches(2.0))
    tf_meta = meta_box.text_frame
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = "Prepared For: Bluestock Fintech Pvt. Ltd.\nProject Type: Capstone Presentation\nPresented By: Intern / Data Analyst\nDate: June 2026"
    p_meta.font.name = 'Arial'
    p_meta.font.size = Pt(13)
    p_meta.font.color.rgb = slate_gray

    # --- SLIDE 2: PROBLEM STATEMENT ---
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "The Industry Problem Landscape")
    add_bullet_content(slide, [
        "Retail investors and advisory firms struggle due to several gaps:",
        "Data Fragmentation",
        "• NAV, AUM, portfolio and inflows are split in static forms (TXT/PDF).",
        "Performance Benchmarking Gaps",
        "• Hard to compare multi-AMC funds on a risk-adjusted basis.",
        "Outperformance Blind Spot",
        "• Investors lack easy tracking of fund returns vs. benchmark indices.",
        "Static Reports",
        "• Reliance on static spreadsheets that take days to compile manually."
    ], width=Inches(8.5))

    # --- SLIDE 3: OBJECTIVES ---
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Project Objectives & Outcomes")
    add_bullet_content(slide, [
        "Goal: Build an institution-grade financial analytics pipeline.",
        "Data Engineering Pipeline",
        "• Automate ingestion from APIs and structure a SQLite DB schema.",
        "Quantitative Risk Metrics",
        "• Calculate CAGR, Sharpe, Sortino, Alpha, Beta, VaR, and CVaR.",
        "Interactive Insights Dashboard",
        "• Streamlit multi-page web app with real-time slicing and filters.",
        "Advanced Behavioral Analytics",
        "• Perform investor cohorts and predict transaction lapse risks."
    ], width=Inches(8.5))

    # --- SLIDE 4: SYSTEM ARCHITECTURE ---
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Platform System Architecture")
    add_bullet_content(slide, [
        "Classic ETL (Extract, Transform, Load) Pipeline:",
        "Extract Layer (Ingestion)",
        "• Scrape mfapi.in API and merge with 10 provided AMFI datasets.",
        "Transform Layer (Processing)",
        "• Standardize dates, handle holidays via forward-fill in Pandas.",
        "Load Layer (SQLite DB)",
        "• Load cleaned datasets into an 8-table relational Star Schema.",
        "Expose Layer (BI Dashboard)",
        "• Power the Streamlit Web Application with fast SQLite connections."
    ], width=Inches(8.5))

    # --- SLIDE 5: DATA INGESTION & QUALITY ---
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Data Ingestion & Quality Reports")
    add_bullet_content(slide, [
        "Programmatic Inspection & Verification of 10 CSVs:",
        "Fund Master (40 schemes) & NAV history (~64k rows) validated.",
        "Cross-verification checked: All 40 AMFI codes successfully mapped.",
        "Missing Values: Handled null values (e.g. yoy_growth_pct) in raw files.",
        "Live Scrapers: Coded live_nav_fetch.py to request real-time historical NAV tables from API endpoints."
    ], width=Inches(8.5))

    # --- SLIDE 6: DATA CLEANING & REINDEXING ---
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Data Cleaning & Reindexing Details")
    add_bullet_content(slide, [
        "Ensuring clean, analytical datasets:",
        "Holiday & Weekend NAV Gaps",
        "• Reindexed dates daily and forward-filled missing NAV rates.",
        "Transaction Formatting",
        "• Standardized types (SIP, Lumpsum, Redemption) and KYC flags.",
        "Constraint Validations",
        "• Verified positive transaction amounts and expense ratio ranges.",
        "Database Execution",
        "• Populated dim_date lookup tables, computing daily return rates."
    ], width=Inches(8.5))

    # --- SLIDE 7: SQL ANALYTICS ---
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "SQL Analytical Queries")
    add_bullet_content(slide, [
        "10 core business queries executing on our Star Schema database:",
        "Top 5 Funds by AUM: Mirae Asset Emerging Bluechip leads at Rs. 49,046 Cr.",
        "Monthly average NAV for SBI Bluechip computed over 53 months.",
        "YoY SIP Inflows Growth calculated dynamically using SQL window LAG.",
        "Investor transactions by payment mode and state volume summaries.",
        "Low expense ratio direct funds (expense < 1.0%) sorted.",
        "Compliance volumes verified (KYC verified vs pending split)."
    ], width=Inches(8.5))

    # --- SLIDE 8: EDA VISUALIZATIONS ---
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "EDA: NAV Trends & Market Growth")
    add_bullet_content(slide, [
        "NAV Historical Movement:",
        "• Solid post-2023 rally seen across Large Cap schemes.",
        "AMC AUM comparison:",
        "• SBI MF leads, followed by ICICI Pru and HDFC MF.",
        "Category Heatmap:",
        "• Net inflows heavily concentrate in Liquid cash funds."
    ], width=Inches(4.5))
    
    # Embed NAV Trends chart
    nav_trends_path = os.path.join(charts_dir, "nav_trends.png")
    if os.path.exists(nav_trends_path):
        slide.shapes.add_picture(nav_trends_path, Inches(5.0), Inches(1.3), Inches(4.5), Inches(2.6))
        
    aum_path = os.path.join(charts_dir, "aum_growth.png")
    if os.path.exists(aum_path):
        slide.shapes.add_picture(aum_path, Inches(5.0), Inches(4.0), Inches(4.5), Inches(2.6))

    # --- SLIDE 9: INVESTOR DEMOGRAPHICS ---
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "EDA: Investor Segmentation")
    add_bullet_content(slide, [
        "Age Group Splits:",
        "• 26-35 age bracket dominates active count (36%).",
        "Geographic analysis:",
        "• Punjab, Tamil Nadu and Rajasthan have top SIP volumes.",
        "Folio Expansion:",
        "• Total folios doubled from 13.26 Cr (2022) to 26.12 Cr (2025)."
    ], width=Inches(4.5))
    
    demo_path = os.path.join(charts_dir, "investor_demographics.png")
    if os.path.exists(demo_path):
        slide.shapes.add_picture(demo_path, Inches(5.0), Inches(1.3), Inches(4.5), Inches(2.6))
        
    geo_path = os.path.join(charts_dir, "geo_distribution.png")
    if os.path.exists(geo_path):
        slide.shapes.add_picture(geo_path, Inches(5.0), Inches(4.0), Inches(4.5), Inches(2.6))

    # --- SLIDE 10: PERFORMANCE METRICS ---
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Quantitative Fund Performance")
    add_bullet_content(slide, [
        "Key mathematical returns and risk indicators calculated:",
        "CAGR Returns: Calculated for 1-year, 3-year, and 5-year periods.",
        "Sharpe Ratio: Measures risk-adjusted excess returns over 6.5% Rf rate.",
        "Sortino Ratio: Penalizes only downside volatility (negative returns).",
        "Beta & Alpha: Aligned index prices to regress fund performance.",
        "Maximum Drawdowns: Traced worst-case peak-to-trough drop rates.",
        "Scorecard: Weighted ranks to construct a composite score out of 100."
    ], width=Inches(8.5))

    # --- SLIDE 11: ADVANCED RISK ANALYTICS ---
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Advanced Analytics & Risk Modules")
    add_bullet_content(slide, [
        "Value at Risk (VaR 95%) & CVaR",
        "• Daily historical worst-case expected loss metrics for each fund.",
        "HHI Sector Concentration Index",
        "• Quantifies portfolio concentration risk. Financial services average > 25%.",
        "SIP Churn Prediction Model",
        "• Flags investors with transaction gaps > 35 days as 'at-risk'.",
        "Risk Recommender Engine (recommender.py)",
        "• Programmatic lookup matching risk appetites to top Sharpe ratio funds."
    ], width=Inches(8.5))

    # --- SLIDE 12: RECOMMENDATIONS & CONCLUSION ---
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Strategic Business Recommendations")
    add_bullet_content(slide, [
        "Actionable insights derived for Bluestock Fintech:",
        "1. Deploy proactive reminder campaigns for flagged 'at-risk' SIP clients.",
        "2. Establish direct-fund distribution outlets in B30 geographical tiers.",
        "3. Highlight direct low-expense schemes to promote platform loyalty.",
        "4. Embed the automated risk recommender engine into the client facing UI.",
        "THANK YOU — QUESTIONS & DISCUSSION"
    ], width=Inches(8.5))

    # Save presentation
    prs.save(prs_path)
    print("PowerPoint Presentation slides compiled successfully.")

if __name__ == "__main__":
    create_presentation()
